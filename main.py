import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import UnstructuredFileLoader,DirectoryLoader
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from PyPDF2 import PdfReader
from pptx import Presentation
from streamlit_chat import message
from langchain.schema import AIMessage, HumanMessage



os.environ["GOOGLE_API_KEY"] = "AIzaSyBfRzQtaS_d6pDoAx-eU-IqCrfQUBr0_Jo"



def get_pdf_text(uploaded_file):
    if uploaded_file is not None:

        for i in uploaded_file:
            text = ''
            if i.name.split('.')[-1] == 'pdf':
                pdf = PdfReader(i)
                for page in pdf.pages:
                    text+='\n\n'+page.extract_text()+'\n\n'
                
            elif i.name.split('.')[-1] == 'pptx':
                prs = Presentation(i)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape,'text'):
                            text+='\n\n'+shape.text

        return text



def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text=text)
    return chunks


def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")


def get_conversational_chain():

    prompt_template = """
    NOTE:
    1) Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in provided context just say, "Answer is not available in the context", don't provide the wrong answer.

    2) You have access to users local files.

    3)Try to give detailed answers. User might also give you a condition related to marks, answer questions acording to the marks.
    
    4) Display the response in beautiful manner this will help user to understand it properly.
    
    5) May be you sometimes receive Greetings message or applause from the user, response them in gentle manner.

    6) Your identity is, you are "Docsxer"
    \n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro",
                             temperature=0.3)

    prompt = PromptTemplate(template = prompt_template, input_variables = ["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    return chain



def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
    
    new_db = FAISS.load_local("faiss_index", embeddings = embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question,k=20)

    chain = get_conversational_chain()

    
    response = chain.invoke(
        {"input_documents":docs, "question": user_question}
        , return_only_outputs=True)

    return response["output_text"]




def main():
    st.set_page_config("Chat PDF")

    with st.sidebar:
        st.title("Menu:")
        pdf_docs = st.file_uploader("Upload your PDF Files and Click on the Submit & Process Button", accept_multiple_files=True)
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text = get_pdf_text(pdf_docs)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Done")

    text = """Hey there, great to meet you. I'm Docsxer, your personal Question-Answering chatbot.
            My goal is to help you in finding you answers from your study material.
            Provide me related document and start a chit chat for accurate and full scoring answers
            """
    
    message(text,is_user=False)

    if 'messages' not in st.session_state:
        st.session_state.messages = []


    user_question = st.chat_input("Ask a Question from the PDF Files")


    if user_question:
        st.session_state.messages.append(HumanMessage(content=user_question))
        response = user_input(user_question)
        st.session_state.messages.append(AIMessage(content=response))

    messages = st.session_state.get('messages',[])
    
    
    for num,res in enumerate(messages):
        if num%2 == 0:
            message(res.content,is_user=True,key = str(num)+'_user')
        else:
            message(res.content,is_user=False,key = str(num)+'_system')






if __name__ == "__main__":
    main()
